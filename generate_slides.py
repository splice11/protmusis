#!/usr/bin/env python3
"""Generate the Trivia Night deck (16 June, Permanent Representation of
Lithuania to the EU) from the questions drafted in `Klausimai protmūšiui.docx`.

Format: one slide per question, followed by one slide with the answer.
Question illustrations are real photos from `photos/` where available
(see photos/CREDITS.md), with OpenMoji icons (CC BY-SA 4.0, rendered to
`assets/` by `fetch_assets.py`) as fallback. The deck is set in Segoe UI
(and Segoe UI Black for display) — fonts installed by default on Windows,
so nothing has to be embedded and managed Office installs render it as
intended.

Song excerpts (media/audio/) and video clips (media/video/) are embedded
to play inline in PowerPoint, each behind an on-brand poster frame
generated with Pillow.

Usage:  pip install python-pptx Pillow && python3 generate_slides.py
Output: Trivia_Night_2026-06-16.pptx
"""

import os
import tempfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette --
BG = RGBColor(0x0B, 0x0E, 0x15)         # near-black background (every slide)
PANEL = RGBColor(0x17, 0x1C, 0x27)      # option cards / fact panels
PAPER = RGBColor(0xFF, 0xFF, 0xFF)      # headline text
LIGHT = RGBColor(0xDD, 0xE2, 0xEB)      # body text on dark
FOG = RGBColor(0x95, 0x9E, 0xAF)        # secondary text on dark
INK = RGBColor(0x10, 0x13, 0x18)        # text on bright slabs


def mix(a, b, t):
    """Blend colour a toward colour b by factor t (0..1)."""
    return RGBColor(*(round(a[i] + (b[i] - a[i]) * t) for i in range(3)))


# One colour per round, tuned as a harmonious jewel-tone set: EU blue, the
# Lithuanian tricolour (gold, green, red) and a twilight indigo for music.
# Each slab is a deep, slightly desaturated base; each bright is a luminous
# (not neon) accent that reads on the dark slides.
# slab = full-bleed divider background · bright = accent on dark slides ·
# on = text colour that reads on the slab.
THEMES = {
    "R1": {"slab": RGBColor(0x1E, 0x3C, 0x86),
           "bright": RGBColor(0x8C, 0xB0, 0xFF), "on": PAPER},
    "R2": {"slab": RGBColor(0xE7, 0xA5, 0x1C),
           "bright": RGBColor(0xFB, 0xC7, 0x4A), "on": INK},
    "R3": {"slab": RGBColor(0x16, 0x77, 0x55),
           "bright": RGBColor(0x57, 0xD2, 0x9E), "on": PAPER},
    "RB": {"slab": RGBColor(0xC2, 0x3B, 0x3B),
           "bright": RGBColor(0xFF, 0x8A, 0x7B), "on": PAPER},
    # Music round — a twilight indigo: distinct from the others, calm and
    # premium, the colour of a concert at dusk by the sea.
    "RM": {"slab": RGBColor(0x46, 0x3C, 0x8E),
           "bright": RGBColor(0xA9, 0x9D, 0xFF), "on": PAPER},
}
TRICOLOUR = [THEMES["R2"]["slab"], THEMES["R3"]["slab"], THEMES["RB"]["slab"]]
CURRENT = THEMES["R2"]  # active round theme; dividers switch it

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EDGE = 0.55                  # breathing room kept clear at the slide's edges
RIGHT = 13.333 - EDGE       # right limit for photos and embedded media (in)
FONT = "Segoe UI"            # installed by default on Windows — no embedding
DISPLAY = "Segoe UI Black"   # heavy display weight, also a Windows system face

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

FOOTER = "Trivia Night · 16 June 2026 · Permanent Representation of Lithuania to the EU"


# ---------------------------------------------------------------- helpers --
def add_slide(bg=None, notes=None, footer=None, show_footer=True):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg or BG
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    if show_footer:
        text(slide, Inches(0.7), Inches(7.08), Inches(10), Inches(0.3),
             FOOTER, size=9.5, color=footer or mix(FOG, BG, 0.35))
    return slide


def shape(slide, kind, x, y, w, h, fill=None, radius=None):
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    if radius is not None:
        sp.adjustments[0] = radius
    return sp


def text(slide, x, y, w, h, content, size=16, color=LIGHT, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None,
         space_after=None, font=FONT):
    """content: str | list of paragraphs; a paragraph is str | list of
    (text, fmt) run tuples with fmt keys bold/color/size/italic."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(content, str):
        content = [content]
    for i, para in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        if space_after is not None:
            p.space_after = space_after
        if isinstance(para, str):
            para = [(para, {})]
        for run_text, fmt in para:
            r = p.add_run()
            r.text = run_text
            r.font.name = font
            r.font.size = Pt(fmt.get("size", size))
            r.font.bold = fmt.get("bold", bold)
            r.font.italic = fmt.get("italic", False)
            r.font.color.rgb = fmt.get("color", color)
    return tb


def icon_tile(slide, name, cx, cy, d, accent=None):
    """Icon on a sharp square tile, centred at (cx, cy) inches, side d.

    The OpenMoji art needs a light tile to read, but a flat cream tile
    clashes with the cool round colours — so the tile is a near-white
    tinted toward the round (or a given) accent, with a hairline edge so
    it stays defined on any background."""
    accent = accent if accent is not None else CURRENT["bright"]
    fill = mix(PAPER, accent, 0.16)
    cx, cy, d = Inches(cx), Inches(cy), Inches(d)
    sp = shape(slide, MSO_SHAPE.RECTANGLE, cx - d / 2, cy - d / 2, d, d,
               fill=fill)
    sp.line.color.rgb = mix(fill, INK, 0.16)
    sp.line.width = Pt(0.75)
    pic = Emu(int(d * 0.72))
    slide.shapes.add_picture(f"assets/{name}.png", cx - pic / 2, cy - pic / 2,
                             pic, pic)


def giant(slide, txt, color, x, y, w, h, size, align=PP_ALIGN.RIGHT,
          anchor=MSO_ANCHOR.TOP):
    """Oversized display glyph (question numbers, divider numerals)."""
    text(slide, Inches(x), Inches(y), Inches(w), Inches(h), txt, size=size,
         color=color, align=align, anchor=anchor, font=DISPLAY)


def edge_bar(slide):
    """Full-height accent bar on the left edge of content slides."""
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), SLIDE_H,
          fill=CURRENT["bright"])


def chip(slide, letter, x, y, side=0.5):
    """Bright square chip with the option letter."""
    sp = shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
               Inches(side), Inches(side), fill=CURRENT["bright"])
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = letter
    r.font.name = DISPLAY
    r.font.size = Pt(15)
    r.font.color.rgb = BG


def tricolour_bars(slide, cy, w=1.05, h=0.14, gap=0.2):
    """Lithuanian tricolour, centred horizontally at height cy (inches)."""
    x = (13.333 - (3 * w + 2 * gap)) / 2
    for c in TRICOLOUR:
        shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(cy), Inches(w),
              Inches(h), fill=c)
        x += w + gap


def place_photo(slide, name, x0, y0, x1, y1, max_h=None, border=True,
                anchor="right"):
    """Place a photo from photos/ into the box (inches) at its native
    aspect ratio — never cropped. The photo fills the box in whichever
    dimension it can; it is anchored to the right edge (or centred) and
    vertically centred. Returns the placed (x, y, w, h)."""
    path = f"photos/{name}"
    src_w, src_h = Image.open(path).size
    aspect = src_w / src_h
    # keep a clear margin at the right edge so images don't crowd it, and
    # keep the accent border off the top/bottom edges
    x1 = min(x1, RIGHT)
    y0, y1 = max(y0, 0.04), min(y1, 7.46)
    bw, bh = x1 - x0, y1 - y0
    if bw / bh > aspect:
        h, w = bh, bh * aspect
    else:
        w, h = bw, bw / aspect
    if max_h and h > max_h:
        h, w = max_h, max_h * aspect
    x = x1 - w if anchor == "right" else x0 + (bw - w) / 2
    y = y0 + (bh - h) / 2
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w),
                                   Inches(h))
    pic.shadow.inherit = False
    if border:
        pic.line.color.rgb = CURRENT["bright"]
        pic.line.width = Pt(1.25)
    return x, y, w, h


def kicker(slide, left, right=None):
    text(slide, Inches(0.7), Inches(0.55), Inches(8.5), Inches(0.35), left,
         size=13, bold=True, color=CURRENT["bright"])
    if right:
        text(slide, Inches(9.2), Inches(0.55), Inches(3.43), Inches(0.35),
             right, size=13, bold=True, color=FOG, align=PP_ALIGN.RIGHT)


# --------------------------------------------------- inline audio / video --
# Song excerpts and video clips are embedded so they play inside PowerPoint,
# offline — no YouTube link to open. python-pptx's add_movie() handles both
# audio and video; each sits behind an on-brand poster frame we render here
# with Pillow so the slide still looks designed before anyone presses play.
_POSTER_DIR = tempfile.mkdtemp(prefix="trivia_posters_")
_POSTER_FONT = "fonts/Montserrat-Black.ttf"  # baked into the PNG, not the deck


def _poster_font(size):
    return ImageFont.truetype(_POSTER_FONT, size)


def make_poster(key, accent, label, sub=None, src=None, size=(1280, 720)):
    """Render a media poster PNG (dark card, accent play button, label) to a
    temp file and return its path. `src` optional background art (drawn to
    fit, then dimmed); `accent` an RGBColor; baked at the box aspect so the
    poster is never stretched."""
    W, H = size
    acc = (accent[0], accent[1], accent[2])
    img = Image.new("RGB", (W, H), (0x0B, 0x0E, 0x15))
    if src:
        art = Image.open(src).convert("RGBA")
        scale = min(W / art.width, H / art.height) * 0.92
        art = art.resize((max(1, int(art.width * scale)),
                          max(1, int(art.height * scale))), Image.LANCZOS)
        img.paste(art, ((W - art.width) // 2, (H - art.height) // 2), art)
        img = Image.alpha_composite(
            img.convert("RGBA"),
            Image.new("RGBA", (W, H), (0x0B, 0x0E, 0x15, 140))).convert("RGB")
    d = ImageDraw.Draw(img)
    cx, cy, r = W // 2, int(H * 0.40), 86
    d.ellipse([cx - r, cy - r, cx + r, cy + r], fill=acc)
    t = 42
    d.polygon([(cx - t * 0.5, cy - t), (cx - t * 0.5, cy + t),
               (cx + t * 0.9, cy)], fill=(0x0B, 0x0E, 0x15))
    f1 = _poster_font(58)
    d.text((cx - d.textlength(label, font=f1) / 2, cy + r + 26), label,
           font=f1, fill=(255, 255, 255))
    if sub:
        f2 = _poster_font(32)
        d.text((cx - d.textlength(sub, font=f2) / 2, cy + r + 104), sub,
               font=f2, fill=(0x95, 0x9E, 0xAF))
    bw = 9
    d.rectangle([bw // 2, bw // 2, W - bw // 2 - 1, H - bw // 2 - 1],
                outline=acc, width=bw)
    path = os.path.join(_POSTER_DIR, key + ".png")
    img.save(path)
    return path


def _set_fullscreen(slide, shape_id):
    """Tick PowerPoint's “Play Full Screen” for the embedded video whose
    shape id is given (sets fullScrn on its media node in the timeline)."""
    timing = slide.element.find(qn("p:timing"))
    if timing is None:
        return
    for node in timing.iter(qn("p:cMediaNode")):
        tgt = node.find(".//" + qn("p:spTgt"))
        if tgt is not None and tgt.get("spid") == str(shape_id):
            node.set("fullScrn", "1")


def media_box(slide, media_path, poster_path, mime, x, y, w, h,
              fullscreen=False):
    """Inline media player (audio or video) framed in the round accent.
    `fullscreen` ticks Play Full Screen so a video fills the screen on
    click (use for video, not the music excerpts)."""
    pad = 0.06
    shape(slide, MSO_SHAPE.RECTANGLE, Inches(x - pad), Inches(y - pad),
          Inches(w + 2 * pad), Inches(h + 2 * pad), fill=CURRENT["bright"])
    gf = slide.shapes.add_movie(media_path, Inches(x), Inches(y), Inches(w),
                                Inches(h), poster_frame_image=poster_path,
                                mime_type=mime)
    if fullscreen:
        _set_fullscreen(slide, gf.shape_id)
    return gf


def link_line(slide, x, y, w, label, url, color, size=11,
              align=PP_ALIGN.RIGHT):
    """A single clickable line of text (a real hyperlink)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w),
                                  Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.hyperlink.address = url
    return tb


def place_movie_right(slide, movie, y=1.75):
    """Embed a 16:9 video clip on the right of a content slide; video plays
    full screen on click. A hyperlinked title sits below it as a YouTube
    backup. Returns the video's left edge x (inches)."""
    vw = 5.95
    vh = vw * 9 / 16
    vx = RIGHT - vw
    poster = make_poster(movie["key"], CURRENT["bright"],
                         movie.get("label", "PLAY THE CLIP"),
                         sub=movie.get("sub"), src=movie.get("poster_src"))
    media_box(slide, movie["file"], poster, movie.get("mime", "video/mp4"),
              vx, y, vw, vh, fullscreen=True)
    cy = y + vh + 0.13
    if movie.get("url"):
        title = movie.get("title", "Watch on YouTube")
        link_line(slide, vx, cy, vw, f"▶  Backup on YouTube — {title}",
                  movie["url"], CURRENT["bright"], size=11)
    else:
        text(slide, Inches(vx), Inches(cy), Inches(vw), Inches(0.3),
             "▶  Embedded — click to play full screen", size=10, color=FOG,
             align=PP_ALIGN.RIGHT)
    return vx


def fact_panel(slide, fact, x, y, w, h):
    shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
          Inches(w), Inches(h), fill=PANEL)
    shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.14),
          Inches(0.07), Inches(h - 0.28), fill=CURRENT["bright"])
    text(slide, Inches(x + 0.4), Inches(y + 0.3), Inches(w - 0.75),
         Inches(0.3), "DID YOU KNOW?", size=12.5, bold=True,
         color=CURRENT["bright"])
    text(slide, Inches(x + 0.4), Inches(y + 0.72), Inches(w - 0.75),
         Inches(h - 1.0), fact, size=15, color=LIGHT, spacing=1.15)


# ----------------------------------------------------------------- slides --
def title_slide():
    # no footer here: the date and venue are on the slide itself, and the
    # ghost question mark bleeds through the footer's spot
    s = add_slide(show_footer=False, notes=(
        "Host notes: welcome everyone and explain the team draw — slips with "
        "names of Lithuanian rivers and lakes at the door."))
    dim = mix(THEMES["R2"]["bright"], BG, 0.84)
    giant(s, "?", dim, 10.4, 0.0, 2.6, 3.4, size=230)
    giant(s, "?", dim, 0.15, 5.35, 2.3, 2.15, size=190, align=PP_ALIGN.LEFT)
    text(s, Inches(0.8), Inches(2.2), Inches(11.73), Inches(0.45),
         "BRUSSELS · 16 JUNE 2026", size=16, bold=True,
         color=THEMES["R2"]["bright"], align=PP_ALIGN.CENTER)
    text(s, Inches(0.4), Inches(2.7), Inches(12.53), Inches(1.65),
         "TRIVIA NIGHT", size=92, color=PAPER, align=PP_ALIGN.CENTER,
         font=DISPLAY)
    tricolour_bars(s, 4.6)
    text(s, Inches(1.5), Inches(4.95), Inches(10.33), Inches(0.9),
         ["Permanent Representation of Lithuania to the European Union",
          [("Environment Team", {"bold": True, "color": LIGHT})]],
         size=16, color=FOG, align=PP_ALIGN.CENTER, spacing=1.3)


def rules_slide():
    s = add_slide(notes=(
        "Team assignment: everyone draws a slip when entering — the slips "
        "carry names of Lithuanian rivers and lakes."))
    text(s, Inches(0.7), Inches(0.6), Inches(12), Inches(0.85),
         "HOW TONIGHT WORKS", size=38, color=PAPER, font=DISPLAY)
    cards = [
        ("teams", "TEAMS", "R1",
         "Draw a slip at the door — the river or lake on it is your team."),
        ("target", "ROUNDS", "R2",
         "Four rounds plus a bonus round. One point per question unless "
         "marked otherwise."),
        ("memo", "ANSWERS", "R3",
         "Write your answers down — we reveal and score after each round."),
        ("no_phone", "FAIR PLAY", "RB",
         "No phones. Diplomatic immunity does not cover googling."),
    ]
    grid = [(0.7, 1.9), (6.95, 1.9), (0.7, 4.35), (6.95, 4.35)]
    for (gx, gy), (icon, head, theme, body) in zip(grid, cards):
        accent = THEMES[theme]["bright"]
        shape(s, MSO_SHAPE.RECTANGLE, Inches(gx), Inches(gy),
              Inches(5.7), Inches(2.05), fill=PANEL)
        shape(s, MSO_SHAPE.RECTANGLE, Inches(gx), Inches(gy),
              Inches(0.07), Inches(2.05), fill=accent)
        icon_tile(s, icon, gx + 1.0, gy + 1.02, 1.15, accent=accent)
        text(s, Inches(gx + 1.85), Inches(gy + 0.38), Inches(3.65),
             Inches(0.42), head, size=17, bold=True, color=accent)
        text(s, Inches(gx + 1.85), Inches(gy + 0.85), Inches(3.65),
             Inches(1.05), body, size=14.5, color=LIGHT, spacing=1.1)


def divider(theme, round_no, title, sub, icons, num, notes=None, size=54):
    global CURRENT
    CURRENT = THEMES[theme]
    slab, on = CURRENT["slab"], CURRENT["on"]
    s = add_slide(bg=slab, notes=notes, footer=mix(on, slab, 0.5))
    giant(s, num, mix(on, slab, 0.85), 6.8, 2.0, 5.83, 5.5, size=300,
          anchor=MSO_ANCHOR.BOTTOM)
    text(s, Inches(0.9), Inches(1.45), Inches(10), Inches(0.45), round_no,
         size=17, bold=True, color=mix(on, slab, 0.2))
    text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.3),
         title.upper(), size=size, color=on, font=DISPLAY, spacing=1.02)
    text(s, Inches(0.9), Inches(4.5), Inches(10), Inches(0.45), sub,
         size=17, bold=True, color=mix(on, slab, 0.25))
    x = 0.9 + 0.55
    for n in icons:
        icon_tile(s, n, x, 6.0, 1.1)
        x += 1.1 + 0.4


def question_slide(round_label, q_no, q_total, question, icon=None,
                   options=None, hint=None, movie=None, notes=None,
                   points="1 point", q_size=26, photo=None, photo_h=None,
                   photo_border=True):
    s = add_slide(notes=notes)
    edge_bar(s)
    kicker(s, round_label, f"QUESTION {q_no} OF {q_total} · {points.upper()}")
    bright = CURRENT["bright"]
    dim = mix(bright, BG, 0.8)
    if options:
        giant(s, str(q_no), dim, 10.6, 0.45, 2.33, 2.1, size=140)
        text(s, Inches(0.7), Inches(1.25), Inches(9.7), Inches(2.4),
             question, size=q_size, bold=True, color=PAPER, spacing=1.12)
        full_row = any(len(o[1]) > 38 for o in options)
        if full_row:
            y = 4.0 if len(options) == 3 else 3.7
            for letter, opt in options:
                shape(s, MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y),
                      Inches(11.93), Inches(0.78), fill=PANEL)
                chip(s, letter, 0.94, y + 0.16, side=0.46)
                text(s, Inches(1.62), Inches(y), Inches(10.7), Inches(0.78),
                     opt, size=17, bold=True, color=LIGHT,
                     anchor=MSO_ANCHOR.MIDDLE)
                y += 0.94
        else:
            gy0 = 3.55 if len(question) < 130 else 4.15
            grid = [(0.7, gy0), (6.78, gy0), (0.7, gy0 + 1.2),
                    (6.78, gy0 + 1.2)]
            for (gx, gy), (letter, opt) in zip(grid, options):
                shape(s, MSO_SHAPE.RECTANGLE, Inches(gx), Inches(gy),
                      Inches(5.85), Inches(0.95), fill=PANEL)
                chip(s, letter, gx + 0.27, gy + 0.225)
                text(s, Inches(gx + 1.02), Inches(gy), Inches(4.6),
                     Inches(0.95), opt, size=18, bold=True, color=LIGHT,
                     anchor=MSO_ANCHOR.MIDDLE)
    else:
        if movie:
            px = place_movie_right(s, movie, y=2.0)
            qw = px - 0.7 - 0.4
        elif photo:
            # borderless art (logos) gets an inset from the slide edge
            x1 = 13.333 if photo_border else 12.93
            px, _, _, _ = place_photo(s, photo, 8.0, 0.0, x1, 7.5,
                                      max_h=photo_h, border=photo_border)
            qw = px - 0.7 - 0.4
        else:
            giant(s, str(q_no), dim, 10.6, 0.45, 2.33, 1.9, size=120)
            icon_tile(s, icon, 10.85, 4.1, 2.7)
            qw, px = 7.9, 9.0
        text(s, Inches(0.7), Inches(1.55), Inches(qw), Inches(4.4),
             question, size=q_size, bold=True, color=PAPER, spacing=1.15)
    if hint:
        text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.4),
             [[("HINT   ", {"bold": True, "color": bright, "size": 14}),
               (hint, {"size": 14, "color": FOG, "italic": True})]])


def answer_slide(round_label, q_no, answer, icon=None, fact=None,
                 movie=None, notes=None, a_size=38, photo=None, credit=None,
                 photo_h=None, photo_border=True):
    s = add_slide(notes=notes)
    edge_bar(s)
    kicker(s, round_label, f"ANSWER · QUESTION {q_no}")
    px = 12.63
    if movie:
        px = place_movie_right(s, movie, y=1.7)
    elif photo:
        names = photo if isinstance(photo, list) else [photo]
        if len(names) == 1:
            bot = 7.02 if credit else 7.5
            x, y, w, h = place_photo(s, names[0], 7.7, 0.2 if credit else 0.0,
                                     13.333, bot, max_h=photo_h,
                                     border=photo_border)
            px = x
            credit_y = y + h + 0.07
        else:
            x1, _, _, _ = place_photo(s, names[0], 7.9, 0.15, 13.333, 3.70)
            x2, y2, w2, h2 = place_photo(s, names[1], 7.9, 3.80, 13.333, 7.35)
            px = min(x1, x2)
            credit_y = min(y2 + h2 + 0.07, 7.18)
        if credit:
            text(s, Inches(px), Inches(credit_y), Inches(13.333 - px - 0.2),
                 Inches(0.3), credit, size=8.5, color=FOG,
                 align=PP_ALIGN.RIGHT)
    elif icon:
        icon_tile(s, icon, 10.85, 3.55, 3.0)
        px = 9.2
    text_w = px - 0.7 - 0.4
    text(s, Inches(0.7), Inches(1.45), Inches(text_w), Inches(1.7), answer,
         size=a_size, bold=True, color=CURRENT["bright"], spacing=1.05)
    if fact:
        fact_panel(s, fact, 0.7, 3.65, text_w, 2.95)


def music_question_slide(round_label, q_no, q_total, audio, hint=None,
                         notes=None):
    """“Name this song”: an embedded excerpt that plays inline, centred
    behind a TRACK n poster. The answer is revealed on the next slide."""
    s = add_slide(notes=notes)
    edge_bar(s)
    kicker(s, round_label, f"QUESTION {q_no} OF {q_total} · 1 POINT")
    text(s, Inches(0.7), Inches(1.4), Inches(11.93), Inches(0.9),
         "NAME THIS SONG", size=46, color=PAPER, font=DISPLAY,
         align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(2.4), Inches(11.93), Inches(0.4),
         [[("Bonus point if you can name the artist",
            {"italic": True})]], size=16, color=FOG, align=PP_ALIGN.CENTER)
    poster = make_poster(f"track{q_no}", CURRENT["bright"], f"TRACK {q_no}",
                         sub="PRESS PLAY")
    vw = 6.4
    vh = vw * 9 / 16
    media_box(s, audio, poster, "audio/mpeg", (13.333 - vw) / 2, 2.95,
              vw, vh)
    if hint:
        text(s, Inches(0.7), Inches(6.6), Inches(11.93), Inches(0.4),
             [[("HINT   ", {"bold": True, "color": CURRENT["bright"],
                            "size": 14}),
               (hint, {"size": 14, "color": FOG, "italic": True})]],
             align=PP_ALIGN.CENTER)


def closing_slide():
    global CURRENT
    CURRENT = THEMES["R2"]
    s = add_slide(notes="Count the points and announce the winning team.")
    place_photo(s, "confetti_popper.jpg", 3.9, 0.55, 9.43, 3.3,
                anchor="center")
    text(s, Inches(0.7), Inches(3.55), Inches(11.93), Inches(2.0),
         ["THANK YOU", "FOR PLAYING"], size=54, color=PAPER,
         align=PP_ALIGN.CENTER, font=DISPLAY, spacing=1.04)
    text(s, Inches(1.5), Inches(5.75), Inches(10.33), Inches(0.5),
         "Time to count the points", size=19, bold=True,
         color=THEMES["R2"]["bright"], align=PP_ALIGN.CENTER)
    tricolour_bars(s, 6.5)


# ------------------------------------------------------------------ deck ---
title_slide()
rules_slide()

# ===== ROUND 1 ==============================================================
R1 = "ROUND 1 · EUROPE & FUN FACTS"
divider("R1", "ROUND 1", "Europe & Fun Facts",
        "Multiple choice · 6 questions · 1 point each",
        ["globe", "euro_note", "jeans"], num="1")

question_slide(
    R1, 1, 6, "Which European island changes sovereignty every six months?",
    "globe",
    options=[("A", "Heligoland"), ("B", "Pheasant Island"), ("C", "Jersey"),
             ("D", "Bornholm")])
answer_slide(
    R1, 1, "B — Pheasant Island",
    photo="pheasant_island.jpg", credit="Photo: Wikimedia Commons",
    fact="Under the 1659 Treaty of the Pyrenees, France and Spain share "
         "this tiny uninhabited island in the Bidasoa river. Administration "
         "alternates between the two countries every six months — the "
         "world's oldest condominium.")

question_slide(
    R1, 2, 6, "Which EU symbol was publicly unveiled about 40 years ago?",
    "star",
    options=[("A", "Euro coins"), ("B", "European flag"),
             ("C", "European anthem"), ("D", "Schengen passport")])
answer_slide(
    R1, 2, "B — The European flag",
    photo="berlaymont_flag_1986.jpg",
    credit="Photo: © European Communities, 1986",
    fact="In May 1986 the European flag was raised for the first time "
         "outside the Berlaymont building — twelve gold stars in a circle, "
         "a symbol of unity, unchanged ever since.")

question_slide(
    R1, 3, 6, "What major international award did the European Union "
              "receive in 2012?",
    "medal",
    options=[("A", "Sakharov Prize"), ("B", "Nobel Peace Prize"),
             ("C", "Charlemagne Prize"), ("D", "Right Livelihood Award")])
answer_slide(
    R1, 3, "B — The Nobel Peace Prize",
    photo="nobel_peace_prize_2012.jpg", photo_h=2.2, credit="Photo: Reuters",
    fact="Awarded to the EU in 2012 for over six decades of advancing "
         "peace, reconciliation, democracy and human rights in Europe.")

question_slide(
    R1, 4, 6, "Which euro banknote was nicknamed “Bin Laden”?",
    "euro_note",
    options=[("A", "€50"), ("B", "€100"), ("C", "€200"), ("D", "€500")],
    notes="Source: https://www.theguardian.com/business/2016/may/04/"
          "500-euro-banknote-could-be-scrapped-crime")
answer_slide(
    R1, 4, "D — The €500 note", photo="euro_500_note.jpg", photo_h=2.6,
    fact="The ECB stopped issuing the €500 note in 2019 over its popularity "
         "with money launderers — and because, like its namesake, everyone "
         "knew what it looked like but almost no one had ever seen one.")

question_slide(
    R1, 5, 6, "In 1984, former Nigerian minister Umaru Dikko was kidnapped "
              "in London. His captors planned to smuggle him out of the UK "
              "in a diplomatic crate, relying on the Vienna Convention rule "
              "that diplomatic bags cannot be opened or detained.\n\nWhat "
              "went wrong?",
    "crate", q_size=21,
    options=[("A", "The crate was improperly labelled as diplomatic "
                   "baggage"),
             ("B", "Dikko did not fit inside the crate"),
             ("C", "The aircraft meant to transport him never arrived")],
    notes="Background video: https://www.youtube.com/watch?v=N83Idy9IOmU")
answer_slide(
    R1, 5, "A — The crate was improperly labelled", a_size=32,
    movie={"key": "dikko", "file": "media/video/umaru_dikko_kidnap.mp4",
           "poster_src": "photos/dikko_crate_stansted_1984.jpg",
           "label": "PLAY THE STORY", "title": "The Umaru Dikko affair",
           "url": "https://www.youtube.com/watch?v=N83Idy9IOmU"},
    fact="With no official diplomatic markings, customs officers at "
         "Stansted were entitled to open the crate — and found Dikko "
         "unconscious, accompanied by an anaesthetist.",
    notes="Embedded clip plays in PowerPoint (media/video/"
          "umaru_dikko_kidnap.mp4). Backup: "
          "https://www.youtube.com/watch?v=N83Idy9IOmU")

question_slide(
    R1, 6, 6, "One of the inventors of modern blue jeans was born in "
              "present-day Latvia. He developed a way to make work trousers "
              "much more durable and partnered with Levi Strauss to patent "
              "the idea.\n\nWhat simple innovation made the trousers "
              "famous?",
    "jeans", q_size=22,
    options=[("A", "A zipper"), ("B", "Metal rivets"),
             ("C", "Waterproof fabric"), ("D", "Belt loops")])
answer_slide(
    R1, 6, "B — Metal rivets", photo="jeans_rivets.jpg",
    fact="Jacob Davis, a tailor born in Riga in 1831, reinforced the stress "
         "points of work trousers with copper rivets. He and Levi Strauss "
         "patented the idea in 1873 — and blue jeans were born.")

# ===== ROUND 2 ==============================================================
R2 = "ROUND 2 · LITHUANIA"
divider("R2", "ROUND 2", "Lithuania", "8 questions · 1 point each",
        ["microphone", "mountain", "soup"], num="2", size=76)

question_slide(
    R2, 1, 8, "Every year Lithuania plans to win Eurovision — although "
              "officially we never have.\n\nIn 2006, however, Lithuania "
              "tried a different strategy: entering a song that simply "
              "declared victory before the contest was even over.\n\nWhat "
              "was the title of the song?",
    photo="eurovision_logo_white.png", photo_h=1.35, photo_border=False,
    q_size=22,
    notes="Keep the question on the Eurovision logo only — the clip is on "
          "the answer slide, so playing it here would give the answer away.")
answer_slide(
    R2, 1, "“We Are The Winners”",
    movie={"key": "eurovision", "file": "media/video/"
           "eurovision_we_are_the_winners.mp4",
           "poster_src": "photos/lt_united_eurovision_2006.jpg",
           "label": "PLAY THE CLIP", "title": "LT United — We Are The Winners",
           "url": "https://www.youtube.com/watch?v=DBAdOlQPbwg"},
    fact="LT United, Eurovision 2006. Honourable mention: Lithuanian-German "
         "singer Lena Valaitis took second place in 1981 — we keep looking "
         "for that victory by all possible means.",
    notes="Play the LT United clip here, on the answer. Decide which part to "
          "show — trim in PowerPoint (Playback ▸ Trim Video). Mention Lena "
          "Valaitis as a fun fact, not a question. Backup: "
          "https://www.youtube.com/watch?v=DBAdOlQPbwg")

question_slide(
    R2, 2, 8, "Two Lithuanian mountaineers carried a symbolic national "
              "object to the summit of Mount Everest and scattered it from "
              "the world's highest peak.\n\nWhat was it?",
    photo="mount_everest.jpg",
    notes="Tell the full story and the exact years out loud; keep the slide "
          "short.")
answer_slide(
    R2, 2, "Amber",
    photo=["vitkauskas_everest_1995.jpg", "baltic_amber.jpg"],
    credit="Top photo: V. Vitkauskas archive",
    fact="Vladas Vitkauskas (1995) and Saulius Vilius (2003) both carried "
         "Baltic amber to the top of the world and scattered it from the "
         "summit.")

question_slide(
    R2, 3, 8, "When Lithuanians learn about the interwar period, one "
              "statistic is proudly repeated in schools: in 1938, Lithuania "
              "ranked second in Europe in butter and bacon production — "
              "despite being one of the poorest countries on the "
              "continent.\n\nWhich country ranked first?",
    photo="butter.jpg", q_size=22,
    hint="Think of a recent Council Presidency.")
answer_slide(
    R2, 3, "Denmark", photo="interwar_lithuania_1930s.jpg",
    fact="A statistic still celebrated in Lithuanian classrooms. Denmark, "
         "for its part, no longer wants to be the EU's bacon factory — as "
         "Politico once put it.",
    notes="Politico: “Denmark no longer wants to be EU bacon factory”.")

question_slide(
    R2, 4, 8, "Across Europe people celebrate Midsummer with bonfires.\n\n"
              "In Lithuania, one tradition involves searching for a "
              "mythical flower that supposedly blooms only on that "
              "night.\n\nWhich flower?",
    photo="midsummer_bonfire.jpg")
answer_slide(
    R2, 4, "The fern flower",
    photo="fern_fiddleheads.jpg", credit="Photo: Wikimedia Commons",
    fact="According to tradition, the fern blooms only on Midsummer night — "
         "whoever finds it gains happiness and wisdom. Botanists remain "
         "unconvinced.")

question_slide(
    R2, 5, 8, "1918 was a big year for Lithuania — also for "
              "democracy.\n\nLithuanian women gained something that many "
              "women in Western Europe had to wait decades for.\n\nWhat "
              "was it?",
    photo="signatories_1918.jpg",
    notes="Alternative version: show the centenary stamp (with some details "
          "removed) and ask what occasion it marks. Stamp: "
          "https://manoteises.lt/straipsnis/isleidziamas-pasto-zenklas-"
          "moteru-balsavimo-simtmeciui-lietuvoje-pamineti/")
answer_slide(
    R2, 5, "Voting rights",
    photo="suffrage_stamp_2018.jpg",
    credit="Stamp: Lietuvos paštas, 2018 · design J. Dadonas",
    fact="Lithuanian women gained the vote on 2 November 1918, when the "
         "provisional constitution enshrined equal suffrage — ahead of "
         "much of Western Europe. A commemorative stamp marked the "
         "centenary.")

question_slide(
    R2, 6, 8, "A 16th-century Lithuanian nobleman travelled through Egypt "
              "and bought several unusual souvenirs.\n\nDuring a storm at "
              "sea, terrified sailors blamed the bad weather on them and "
              "threw them overboard.\n\nWhat were the souvenirs?",
    photo="aswan_egypt_nile.jpg", q_size=24)
answer_slide(
    R2, 6, "Egyptian mummies",
    photo="radvila_the_orphan.jpg",
    credit="Anonymous portrait, c. 1590 (public domain)",
    fact="Mikalojus Kristupas Radvila the Orphan brought two mummies back "
         "from his pilgrimage. When storms battered the ship, the crew "
         "blamed the cargo — and overboard they went. His travel diary "
         "became a European bestseller.")

question_slide(
    R2, 7, 8, "In a humorous Lithuanian promotional video from the early "
              "2000s, traditional dishes such as cepelinai are discussed "
              "with the idea that they might one day become popular across "
              "Europe as street food.\n\nThe video reflects a moment when "
              "Lithuania was preparing for a major change. What event was "
              "approaching?",
    q_size=22,
    movie={"key": "spirgi", "file": "media/video/eu_accession_spirgi.mp4",
           "poster_src": "photos/cepelinai.jpg", "label": "PLAY THE CLIP",
           "title": "Spirgi, spirgi — EU accession promo",
           "url": "https://www.youtube.com/watch?v=YgvHcenDYcU"},
    notes="Embedded clip plays in PowerPoint (media/video/"
          "eu_accession_spirgi.mp4) — add English subtitles before the "
          "night. Consider whether to keep “early 2000s” in the wording. "
          "Backup: https://www.youtube.com/watch?v=YgvHcenDYcU")
answer_slide(
    R2, 7, "Joining the European Union", a_size=34,
    photo="eu_enlargement_2004_map.png", credit="Map: Wikimedia Commons",
    fact="Lithuania joined the EU on 1 May 2004, in the largest enlargement "
         "in the Union's history — ten countries at once. The cepelinai "
         "street-food revolution is still pending.")

question_slide(
    R2, 8, 8, "Pink is strongly associated with Lithuania because of a "
              "traditional summer dish.\n\nWhich dish?",
    "blossom",
    options=[("A", "Cepelinai"), ("B", "Šakotis"), ("C", "Šaltibarščiai"),
             ("D", "Kibinai")])
answer_slide(
    R2, 8, "C — Šaltibarščiai", photo="saltibarsciai.jpg",
    fact="The electric-pink cold beet soup is Lithuania's unofficial "
         "summer flag — best served with hot potatoes and a sunny terrace.")

# ===== ROUND 3 ==============================================================
R3 = "ROUND 3 · ENVIRONMENT & NATURE"
# NB: no beaver or headphones on the divider — both are answers in this round
divider("R3", "ROUND 3", "Environment & Nature",
        "4 questions · 1 point each",
        ["tree", "fern", "deer"], num="3", size=50)

question_slide(
    R3, 1, 4, "In 2013 Metallica became the first band to perform on all "
              "seven continents.\n\nDuring their Antarctic concert, the "
              "audience used 120 of what?",
    photo="metallica_freeze_em_all.jpg")
answer_slide(
    R3, 1, "Headphones",
    photo="metallica_antarctica_2013.jpg",
    credit="Photo: Metallica (freezeemall.com)",
    fact="To comply with Antarctic environmental rules, the band played "
         "without amplifiers — the audience of about 120 listened through "
         "headphones. The concert was fittingly called “Freeze 'Em All”.")

question_slide(
    R3, 2, 4, "Three famous scientists — Birutė Galdikas, Dian Fossey and "
              "Jane Goodall — devoted their careers to studying these "
              "animals in the wild.\n\nCollectively, they became known as "
              "the “Trimates”.\n\nWhat group of animals are these?",
    photo="trimates_goodall_fossey_galdikas.jpg", q_size=24)
answer_slide(
    R3, 2, "Great apes",
    photo="great_apes_collage.jpg",
    credit="Collage: The New York Times",
    fact="Goodall (chimpanzees), Fossey (gorillas) and Galdikas "
         "(orangutans) were recruited by Louis Leakey — hence also "
         "“Leakey's Angels”. Galdikas, of Lithuanian descent, still works "
         "in Borneo. “Orangutan” means “person of the forest”.")

question_slide(
    R3, 3, 4, "Which EU law unexpectedly entered the headlines after the "
              "deaths of four US soldiers, whose vehicle sank in a military "
              "training area near the Lithuanian–Belarusian border?",
    photo="pabrade_recovery_2025.jpg")
answer_slide(
    R3, 3, "The Nature Restoration Law", a_size=34,
    photo="cepkeliai_marsh.jpg", credit="Photo: Wikimedia Commons",
    fact="Commentary on the 2025 tragedy near Pabradė linked the swampy "
         "terrain to wetlands restored under EU environmental policy — "
         "putting the Nature Restoration Law unexpectedly in the news.")

question_slide(
    R3, 4, 4, "This animal is often called an “ecosystem engineer” because "
              "it creates wetlands that benefit countless other "
              "species.\n\nWhat animal is it?",
    photo="pond_ecosystem.jpg")
answer_slide(
    R3, 4, "The beaver", photo="beaver_at_dam.jpg",
    fact="Beaver dams create wetlands that store water, filter pollution "
         "and host countless species. Lithuania's beaver population has "
         "grown from near extinction to one of the densest in Europe.")

# ===== ROUND 4 — MUSIC ======================================================
RM = "ROUND 4 · MUSIC"
divider("RM", "ROUND 4", "Name That Tune",
        "5 songs · 1 point each · press play",
        ["microphone", "guitar", "headphones"], num="4", size=64,
        notes="Each slide plays an embedded excerpt — just click the play "
              "button. Teams write down the song (bonus point for the "
              "artist). The twist: every track is about the sea — a nod to "
              "Čiurlionis's symphonic poem “Jūra”. You can reveal that link "
              "at the end of the round for an extra point.")

music_question_slide(
    RM, 1, 5, "media/audio/how_far_ill_go.mp3",
    hint="From a 2016 Disney film about a Pacific island voyager.")
answer_slide(
    RM, 1, "“How Far I'll Go”", a_size=40,
    photo="song_how_far_ill_go.png", credit="Cover: Walt Disney Records",
    fact="Sung by Auli'i Cravalho in Disney's “Moana” (2016) — the ocean "
         "keeps calling her past the reef. Theme of the round: the sea.")

music_question_slide(
    RM, 2, 5, "media/audio/my_heart_will_go_on.mp3",
    hint="The love theme from the highest-grossing film of the 1990s.")
answer_slide(
    RM, 2, "“My Heart Will Go On”", a_size=38,
    photo="song_my_heart_will_go_on.png",
    credit="Single cover: Columbia / Sony",
    fact="Céline Dion's theme from “Titanic” (1997) — an ocean liner, and "
         "another track all at sea.")

music_question_slide(
    RM, 3, 5, "media/audio/jura_happyendless.mp3",
    hint="A Lithuanian act — and the title is the Lithuanian word for "
         "“the sea”.")
answer_slide(
    RM, 3, "“Jūra” — HappyEndless", a_size=36,
    photo="song_jura_happyendless.jpg",
    fact="“Jūra” is Lithuanian for “the sea” — the most direct clue to the "
         "round's hidden theme.")

music_question_slide(
    RM, 4, 5, "media/audio/yellow_submarine.mp3",
    hint="A 1966 sing-along by the most famous band from Liverpool.")
answer_slide(
    RM, 4, "“Yellow Submarine”", a_size=40,
    photo="song_yellow_submarine.jpg",
    credit="Album art: Apple / EMI",
    fact="The Beatles, 1966 — “we all live in a yellow submarine”, somewhere "
         "beneath the sea.")

music_question_slide(
    RM, 5, 5, "media/audio/ocean_eyes.mp3",
    hint="The 2015 breakout track, written when the singer was 13.")
answer_slide(
    RM, 5, "“Ocean Eyes”", a_size=42,
    photo="song_ocean_eyes.png", credit="Single cover: Darkroom / Interscope",
    fact="Billie Eilish's debut, 2015 — and the fifth and final clue: every "
         "song in this round is about the sea.")

# ===== BONUS ROUND ==========================================================
RB = "BONUS · THE BRUSSELS BUBBLE"
divider("RB", "BONUS ROUND", "The Brussels Bubble",
        "Insider questions · for those who were in the room",
        ["recycle", "speech", "ship"], num="B", size=50,
        notes="Possible extra questions: the CBAM question about nails "
              "(still needs wording — a funny one); The Schumann Show about "
              "the institutions: https://www.youtube.com/watch?v=D0fBh-0Eiy0")

question_slide(
    RB, 1, 3, "During negotiations on which file did Violeta Dragu "
              "(Romania) say:\n\n“I have never seen Lithuania being so "
              "active during the negotiations”?",
    "speech",
    notes="Can be made multiple choice to make it easier — or harder, by "
          "also asking what the problem was.")
answer_slide(
    RB, 1, "The Waste Framework Directive", a_size=34,
    photo="waste_sorting_bins.png",
    credit="Illustration: EPRS briefing, European Parliament")

question_slide(
    RB, 2, 3, "During a Coreper I meeting, the Council was split on a "
              "legislative file, with both sides carefully counting "
              "votes.\n\nThe Deputy Permanent Representative of Lithuania "
              "intervened and said… what?\n\nAnd which file was being "
              "discussed?",
    "abacus", q_size=23, points="2 points",
    notes="Optional hint: give the exact date — the 14 June Coreper I (or "
          "the earlier one).")
answer_slide(
    RB, 2, "“Scrutiny reservation” — on the Nature Restoration Law",
    a_size=30, photo="europa_building_room.jpg",
    fact="A scrutiny reservation lets a member state hold its position "
         "while procedures are completed back home — a small phrase that "
         "can pause a finely balanced vote. Pictured: the Europa building, "
         "where Coreper meets.")

question_slide(
    RB, 3, 3, "1985: Lithuanian, Latvian and Estonian dissidents organised "
              "the Peace and Freedom Cruise aboard the ship Baltic Star — "
              "sailing from Stockholm along the edge of Soviet territorial "
              "waters and ending with a widely reported demonstration in "
              "Helsinki.\n\nThe Soviet Union tried to derail the voyage. "
              "What did it do?",
    photo="baltic_star_birger_jarl.jpg", q_size=22,
    notes="Background: per Tomas Venclova, Aleksandras Štromas was the "
          "cruise's spiritus movens; the Baltic Tribunal also took place in "
          "1985. Keep names off the slide. Ship photo: Wikimedia Commons "
          "(M/S Birger Jarl, sailed as Baltic Star).")
answer_slide(
    RB, 3, "A rumour that a bomb was on board", a_size=32,
    photo="baltic_cruise_cartoon_1985.jpg",
    credit="Cartoon, 1985 — exile press, via X",
    fact="The rumour failed: the cruise went ahead, and the Helsinki "
         "demonstration was reported across the Scandinavian and wider "
         "European press — a loud reminder that the Baltic states had not "
         "been forgotten.")

closing_slide()


OUT = "Trivia_Night_2026-06-16.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides "
      "(Segoe UI · audio & video embedded)")
